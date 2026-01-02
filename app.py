from flask import Flask, render_template, request, jsonify, send_file
import json
import os
from datetime import datetime
from chatbot import ChatBot
from werkzeug.utils import secure_filename
import tempfile

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'docx', 'pptx', 'ppt', 'doc'}

# Initialize the chatbot
chatbot = ChatBot()

# Create necessary directories (works on both local and Render)
os.makedirs('uploads', exist_ok=True)
os.makedirs('templates', exist_ok=True)
os.makedirs('static', exist_ok=True)
os.makedirs('models', exist_ok=True)
os.makedirs('data', exist_ok=True)

@app.route('/')
def index():
    return render_template('home.html')

@app.route('/chat')
def chat_page():
    return render_template('chat.html')

@app.route('/converter')
def converter_page():
    return render_template('converter.html')

@app.route('/summarizer')
def summarizer_page():
    return render_template('summarizer.html')

@app.route('/image-generator')
def image_generator_page():
    return render_template('image_generator.html')


@app.route('/chat', methods=['POST'])
def chat():
    try:
        data = request.get_json()
        user_message = data.get('message', '')
        
        if not user_message:
            return jsonify({'error': 'No message provided'}), 400
        
        # Get response from chatbot
        response = chatbot.get_response(user_message)
        
        return jsonify({
            'response': response,
            'timestamp': datetime.now().isoformat()
        })
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/train', methods=['POST'])
def train():
    try:
        data = request.get_json()
        training_data = data.get('training_data', [])
        
        if not training_data:
            return jsonify({'error': 'No training data provided'}), 400
        
        # Train the chatbot with new data
        result = chatbot.train_model(training_data)
        
        return jsonify({
            'message': 'Training completed successfully',
            'trained_samples': result['trained_samples'],
            'timestamp': datetime.now().isoformat()
        })
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/add_data', methods=['POST'])
def add_training_data():
    try:
        data = request.get_json()
        question = data.get('question', '').strip()
        answer = data.get('answer', '').strip()
        
        if not question or not answer:
            return jsonify({'error': 'Both question and answer are required'}), 400
        
        # Add to training data
        success = chatbot.add_training_data(question, answer)
        
        if success:
            return jsonify({
                'message': 'Training data added successfully',
                'timestamp': datetime.now().isoformat()
            })
        else:
            return jsonify({'error': 'Failed to add training data'}), 500
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/get_training_data', methods=['GET'])
def get_training_data():
    try:
        training_data = chatbot.get_training_data()
        return jsonify({
            'training_data': training_data,
            'count': len(training_data)
        })
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/delete_data', methods=['POST'])
def delete_training_data():
    try:
        data = request.get_json()
        index = data.get('index')
        
        if index is None:
            return jsonify({'error': 'Index is required'}), 400
        
        success = chatbot.delete_training_data(index)
        
        if success:
            return jsonify({
                'message': 'Training data deleted successfully',
                'timestamp': datetime.now().isoformat()
            })
        else:
            return jsonify({'error': 'Failed to delete training data'}), 500
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/retrain', methods=['POST'])
def retrain():
    try:
        # Retrain the model with all current training data
        result = chatbot.retrain_model()
        
        return jsonify({
            'message': 'Model retrained successfully',
            'trained_samples': result['trained_samples'],
            'timestamp': datetime.now().isoformat()
        })
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/conversation/summary', methods=['GET'])
def get_conversation_summary():
    try:
        summary = chatbot.get_conversation_summary()
        return jsonify({
            'summary': summary,
            'topics': chatbot.conversation_topics,
            'message_count': len([entry for entry in chatbot.conversation_history if 'user' in entry])
        })
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/conversation/clear', methods=['POST'])
def clear_conversation():
    try:
        message = chatbot.clear_conversation()
        return jsonify({
            'message': message,
            'timestamp': datetime.now().isoformat()
        })
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/conversation/history', methods=['GET'])
def get_conversation_history():
    try:
        return jsonify({
            'history': chatbot.conversation_history[-10:],  # Last 10 messages
            'timestamp': datetime.now().isoformat()
        })
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/personality', methods=['GET'])
def get_personality():
    try:
        personality_info = chatbot.get_personality_info()
        return jsonify(personality_info)
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/personality/adjust', methods=['POST'])
def adjust_personality():
    try:
        data = request.get_json()
        trait = data.get('trait')
        value = data.get('value')
        
        if not trait or value is None:
            return jsonify({'error': 'Trait and value are required'}), 400
        
        result = chatbot.adjust_personality(trait, value)
        return jsonify({
            'message': result,
            'personality': chatbot.get_personality_info()
        })
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

def allowed_file(filename, extension=None):
    if extension:
        return extension.lower() in app.config['ALLOWED_EXTENSIONS']
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def convert_ppt_to_pptx(input_path, output_path):
    """Convert old PowerPoint .ppt to .pptx using LibreOffice or alternative method"""
    import subprocess
    import shutil
    
    # Method 1: Try LibreOffice command line (works on Linux/Mac/Windows if installed)
    try:
        # Try different LibreOffice command names
        libreoffice_cmd = None
        for cmd in ['libreoffice', 'soffice', '/Applications/LibreOffice.app/Contents/MacOS/soffice']:
            if shutil.which(cmd):
                libreoffice_cmd = cmd
                break
        
        if libreoffice_cmd:
            # Use LibreOffice to convert
            cmd = [
                libreoffice_cmd,
                '--headless',
                '--convert-to', 'pptx',
                '--outdir', os.path.dirname(output_path),
                input_path
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode == 0:
                # LibreOffice creates file with same name but .pptx extension
                converted_file = os.path.join(
                    os.path.dirname(output_path),
                    os.path.splitext(os.path.basename(input_path))[0] + '.pptx'
                )
                if os.path.exists(converted_file):
                    if converted_file != output_path:
                        shutil.move(converted_file, output_path)
                    return True
    except Exception as e:
        print(f"LibreOffice conversion failed: {e}")
    
    # Method 2: Try pypandoc if available
    try:
        import pypandoc
        pypandoc.convert_file(input_path, 'pptx', outputfile=output_path)
        if os.path.exists(output_path):
            return True
    except ImportError:
        pass
    except Exception as e:
        print(f"Pypandoc conversion failed: {e}")
    
    # Method 3: Try comtypes on Windows (for .ppt files)
    try:
        import platform
        if platform.system() == 'Windows':
            import comtypes.client
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1
            presentation = powerpoint.Presentations.Open(input_path)
            presentation.SaveAs(output_path, 24)  # 24 = ppSaveAsOpenXMLPresentation
            presentation.Close()
            powerpoint.Quit()
            if os.path.exists(output_path):
                return True
    except ImportError:
        pass
    except Exception as e:
        print(f"Comtypes conversion failed: {e}")
    
    # If all methods fail, raise an error
    raise Exception("Could not convert .ppt to .pptx. Please install LibreOffice or convert the file manually to .pptx format first.")

def convert_doc_to_docx(input_path, output_path):
    """Convert old Word .doc to .docx using LibreOffice or alternative method"""
    import subprocess
    import shutil
    
    # Method 1: Try LibreOffice command line
    try:
        libreoffice_cmd = None
        for cmd in ['libreoffice', 'soffice', '/Applications/LibreOffice.app/Contents/MacOS/soffice']:
            if shutil.which(cmd):
                libreoffice_cmd = cmd
                break
        
        if libreoffice_cmd:
            cmd = [
                libreoffice_cmd,
                '--headless',
                '--convert-to', 'docx',
                '--outdir', os.path.dirname(output_path),
                input_path
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            
            if result.returncode == 0:
                converted_file = os.path.join(
                    os.path.dirname(output_path),
                    os.path.splitext(os.path.basename(input_path))[0] + '.docx'
                )
                if os.path.exists(converted_file):
                    if converted_file != output_path:
                        shutil.move(converted_file, output_path)
                    return True
    except Exception as e:
        print(f"LibreOffice conversion failed: {e}")
    
    # Method 2: Try pypandoc if available
    try:
        import pypandoc
        pypandoc.convert_file(input_path, 'docx', outputfile=output_path)
        if os.path.exists(output_path):
            return True
    except ImportError:
        pass
    except Exception as e:
        print(f"Pypandoc conversion failed: {e}")
    
    # Method 3: Try comtypes on Windows
    try:
        import platform
        if platform.system() == 'Windows':
            import comtypes.client
            word = comtypes.client.CreateObject("Word.Application")
            word.Visible = 0
            doc = word.Documents.Open(input_path)
            doc.SaveAs(output_path, 16)  # 16 = wdFormatDocumentDefault (docx)
            doc.Close()
            word.Quit()
            if os.path.exists(output_path):
                return True
    except ImportError:
        pass
    except Exception as e:
        print(f"Comtypes conversion failed: {e}")
    
    # If all methods fail, raise an error
    raise Exception("Could not convert .doc to .docx. Please install LibreOffice or convert the file manually to .docx format first.")

def convert_pdf_to_word(input_path, output_path):
    """Convert PDF to Word document"""
    from pdf2docx import Converter
    cv = Converter(input_path)
    cv.convert(output_path)
    cv.close()

def convert_word_to_pdf(input_path, output_path):
    """Convert Word document to PDF"""
    from docx import Document
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.units import inch
    
    doc = Document(input_path)
    pdf_doc = SimpleDocTemplate(output_path, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    for para in doc.paragraphs:
        if para.text.strip():
            # Determine style based on paragraph style
            style_name = para.style.name if para.style else 'Normal'
            if 'Heading' in style_name or 'Title' in style_name:
                style = styles['Heading1']
            else:
                style = styles['Normal']
            
            # Clean text and create paragraph
            text = para.text.replace('\n', ' ')
            p = Paragraph(text, style)
            story.append(p)
            story.append(Spacer(1, 0.2*inch))
    
    pdf_doc.build(story)
    
    # Clean up converted file if it was created
    if converted_file and os.path.exists(converted_file):
        try:
            os.remove(converted_file)
        except:
            pass

def convert_pptx_to_pdf(input_path, output_path):
    """Convert PowerPoint to PDF"""
    from pptx import Presentation
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from reportlab.lib.styles import getSampleStyleSheet
    
    # Handle old .ppt files
    actual_input_path = input_path
    converted_file = None
    if input_path.lower().endswith('.ppt'):
        try:
            converted_file = os.path.join(tempfile.gettempdir(), f"temp_pptx_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx")
            convert_ppt_to_pptx(input_path, converted_file)
            if os.path.exists(converted_file):
                actual_input_path = converted_file
        except Exception as e:
            raise Exception(f"Failed to convert .ppt to .pptx: {str(e)}")
    
    try:
        prs = Presentation(actual_input_path)
    except Exception as e:
        raise Exception(f"Failed to open PowerPoint file: {str(e)}")
    
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    styles = getSampleStyleSheet()
    
    for slide_num, slide in enumerate(prs.slides):
        if slide_num > 0:
            c.showPage()
        
        y = height - 50
        c.setFont("Helvetica-Bold", 16)
        c.drawString(50, y, f"Slide {slide_num + 1}")
        y -= 30
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.replace('\n', ' ')
                # Wrap long text
                lines = []
                words = text.split()
                current_line = ""
                for word in words:
                    if len(current_line + word) < 80:
                        current_line += word + " "
                    else:
                        if current_line:
                            lines.append(current_line.strip())
                        current_line = word + " "
                if current_line:
                    lines.append(current_line.strip())
                
                c.setFont("Helvetica", 12)
                for line in lines[:10]:  # Limit to 10 lines per shape
                    if y < 50:
                        c.showPage()
                        y = height - 50
                    c.drawString(70, y, line[:80])
                    y -= 18
        
        if y < 100:
            c.showPage()
    
    c.save()
    
    # Clean up converted file if it was created
    if converted_file and os.path.exists(converted_file):
        try:
            os.remove(converted_file)
        except:
            pass

def convert_pdf_to_pptx(input_path, output_path):
    """Convert PDF to PowerPoint (extract text and create slides)"""
    import PyPDF2
    from pptx import Presentation
    
    prs = Presentation()
    prs.slide_width = 9144000  # 10 inches
    prs.slide_height = 6858000  # 7.5 inches
    
    with open(input_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        for page_num, page in enumerate(pdf_reader.pages):
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            
            text = page.extract_text()
            if text.strip():
                if slide.shapes.title:
                    slide.shapes.title.text = f"Slide {page_num + 1}"
                
                if len(slide.placeholders) > 1 and slide.placeholders[1]:
                    content = slide.placeholders[1]
                    tf = content.text_frame
                    tf.text = ""
                    # Add text in paragraphs
                    text_lines = text.split('\n')[:20]  # Limit to 20 lines
                    for i, line in enumerate(text_lines):
                        if line.strip():
                            if i == 0:
                                tf.text = line[:200]
                            else:
                                p = tf.add_paragraph()
                                p.text = line[:200]
    
    prs.save(output_path)

def convert_word_to_pptx(input_path, output_path):
    """Convert Word to PowerPoint (extract text and create slides)"""
    from docx import Document
    from pptx import Presentation
    
    # Handle old .doc files
    actual_input_path = input_path
    converted_file = None
    if input_path.lower().endswith('.doc'):
        try:
            converted_file = os.path.join(tempfile.gettempdir(), f"temp_docx_{datetime.now().strftime('%Y%m%d%H%M%S')}.docx")
            convert_doc_to_docx(input_path, converted_file)
            if os.path.exists(converted_file):
                actual_input_path = converted_file
        except Exception as e:
            raise Exception(f"Failed to convert .doc to .docx: {str(e)}")
    
    doc = Document(actual_input_path)
    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 6858000
    
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Document Content"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = ""
    
    for para in doc.paragraphs:
        if para.text.strip():
            p = tf.add_paragraph()
            p.text = para.text[:200]
    
    prs.save(output_path)
    
    # Clean up converted file if it was created
    if converted_file and os.path.exists(converted_file):
        try:
            os.remove(converted_file)
        except:
            pass

def convert_pptx_to_word(input_path, output_path):
    """Convert PowerPoint to Word (extract text and create document)"""
    from pptx import Presentation
    from docx import Document
    
    # Handle old .ppt files
    actual_input_path = input_path
    converted_file = None
    if input_path.lower().endswith('.ppt'):
        try:
            converted_file = os.path.join(tempfile.gettempdir(), f"temp_pptx_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx")
            convert_ppt_to_pptx(input_path, converted_file)
            if os.path.exists(converted_file):
                actual_input_path = converted_file
        except Exception as e:
            raise Exception(f"Failed to convert .ppt to .pptx: {str(e)}")
    
    try:
        prs = Presentation(actual_input_path)
    except Exception as e:
        raise Exception(f"Failed to open PowerPoint file: {str(e)}")
    
    doc = Document()
    
    doc.add_heading('Presentation Content', 0)
    
    for slide_num, slide in enumerate(prs.slides):
        doc.add_heading(f'Slide {slide_num + 1}', level=1)
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                doc.add_paragraph(shape.text)
    
    doc.save(output_path)
    
    # Clean up converted file if it was created
    if converted_file and os.path.exists(converted_file):
        try:
            os.remove(converted_file)
        except:
            pass

@app.route('/convert/file', methods=['POST'])
def convert_file():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        from_format = request.form.get('from_format', '').lower()
        to_format = request.form.get('to_format', '').lower()
        
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not from_format or not to_format:
            return jsonify({'error': 'Source and target formats are required'}), 400
        
        if from_format == to_format:
            return jsonify({'error': 'Source and target formats cannot be the same'}), 400
        
        # Validate file extension matches from_format (support both old and new formats)
        file_ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
        
        # Map old formats to new formats for compatibility
        format_mapping = {
            'ppt': 'pptx',  # Old PowerPoint format
            'doc': 'docx'   # Old Word format
        }
        
        # Normalize file extension
        normalized_ext = format_mapping.get(file_ext, file_ext)
        expected_ext = format_mapping.get(from_format, from_format)
        
        # Allow both old and new formats
        valid_extensions = {
            'pptx': ['ppt', 'pptx'],
            'docx': ['doc', 'docx'],
            'pdf': ['pdf']
        }
        
        if normalized_ext not in valid_extensions.get(expected_ext, [expected_ext]):
            return jsonify({'error': f'File extension (.{file_ext}) does not match selected source format ({from_format}). Accepted: {", ".join(valid_extensions.get(expected_ext, [expected_ext]))}'}), 400
        
        if not allowed_file(file.filename, from_format):
            return jsonify({'error': f'Invalid file type. {from_format.upper()} files are not supported.'}), 400
        
        # Save uploaded file temporarily (use tempfile for Render compatibility)
        filename = secure_filename(file.filename)
        # Use tempfile for better compatibility with Render's ephemeral filesystem
        upload_path = os.path.join(tempfile.gettempdir(), f"upload_{datetime.now().strftime('%Y%m%d%H%M%S')}_{filename}")
        file.save(upload_path)
        
        try:
            # Determine output filename and path
            output_filename = filename.rsplit('.', 1)[0] + '.' + to_format
            output_path = os.path.join(tempfile.gettempdir(), output_filename)
            
            # MIME types for different formats
            mime_types = {
                'pdf': 'application/pdf',
                'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            }
            
            # Normalize format (handle old .ppt and .doc to new formats)
            format_mapping = {
                'ppt': 'pptx',
                'doc': 'docx'
            }
            normalized_from = format_mapping.get(from_format, from_format)
            normalized_to = format_mapping.get(to_format, to_format)
            
            # Check if old format is being used - python-pptx only supports .pptx, not .ppt
            if file_ext == 'ppt':
                os.remove(upload_path)
                return jsonify({
                    'error': 'Old PowerPoint format (.ppt) is not supported. The python-pptx library only works with .pptx files. Please convert your .ppt file to .pptx format first using Microsoft PowerPoint or an online converter.'
                }), 400
            
            if file_ext == 'doc':
                os.remove(upload_path)
                return jsonify({
                    'error': 'Old Word format (.doc) is not supported. Please convert your .doc file to .docx format first using Microsoft Word or an online converter.'
                }), 400
            
            # Perform conversion based on format combination (use normalized formats)
            conversion_key = f"{normalized_from}_to_{normalized_to}"
            
            try:
                if conversion_key == 'pdf_to_docx':
                    convert_pdf_to_word(actual_input_path, output_path)
                elif conversion_key == 'docx_to_pdf':
                    convert_word_to_pdf(actual_input_path, output_path)
                elif conversion_key == 'pptx_to_pdf':
                    convert_pptx_to_pdf(actual_input_path, output_path)
                elif conversion_key == 'pdf_to_pptx':
                    convert_pdf_to_pptx(actual_input_path, output_path)
                elif conversion_key == 'docx_to_pptx':
                    convert_word_to_pptx(actual_input_path, output_path)
                elif conversion_key == 'pptx_to_docx':
                    convert_pptx_to_word(actual_input_path, output_path)
                else:
                    # Clean up files
                    os.remove(upload_path)
                    if converted_file_path and os.path.exists(converted_file_path):
                        os.remove(converted_file_path)
                    return jsonify({'error': f'Conversion from {from_format} to {to_format} is not supported'}), 400
            except Exception as conv_error:
                # Clean up files
                if os.path.exists(upload_path):
                    os.remove(upload_path)
                if converted_file_path and os.path.exists(converted_file_path):
                    os.remove(converted_file_path)
                error_msg = str(conv_error)
                if 'not supported' in error_msg.lower() or 'old format' in error_msg.lower():
                    return jsonify({'error': error_msg}), 400
                return jsonify({'error': f'Conversion failed: {error_msg}'}), 500
            
            # Clean up uploaded and converted files
            if os.path.exists(upload_path):
                os.remove(upload_path)
            if converted_file_path and os.path.exists(converted_file_path):
                os.remove(converted_file_path)
            
            # Send the converted file
            return send_file(
                output_path,
                as_attachment=True,
                download_name=output_filename,
                mimetype=mime_types.get(to_format, 'application/octet-stream')
            )
        
        except ImportError as e:
            # Clean up files
            if os.path.exists(upload_path):
                os.remove(upload_path)
            if converted_file_path and os.path.exists(converted_file_path):
                os.remove(converted_file_path)
            return jsonify({'error': f'Required library not installed: {str(e)}'}), 500
        except Exception as e:
            # Clean up files
            if os.path.exists(upload_path):
                os.remove(upload_path)
            if converted_file_path and os.path.exists(converted_file_path):
                os.remove(converted_file_path)
            return jsonify({'error': f'Conversion failed: {str(e)}'}), 500
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

def extract_text_from_pdf(file_path):
    """Extract text from PDF file"""
    import PyPDF2
    text = ""
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        raise Exception(f"Failed to extract text from PDF: {str(e)}")
    return text

def extract_text_from_word(file_path):
    """Extract text from Word document"""
    from docx import Document
    try:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    except Exception as e:
        raise Exception(f"Failed to extract text from Word document: {str(e)}")
    return text

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint presentation"""
    from pptx import Presentation
    
    # Check if it's an old .ppt file
    if file_path.lower().endswith('.ppt'):
        # Try to convert .ppt to .pptx first
        try:
            converted_path = os.path.join(tempfile.gettempdir(), f"temp_pptx_{datetime.now().strftime('%Y%m%d%H%M%S')}.pptx")
            convert_ppt_to_pptx(file_path, converted_path)
            if os.path.exists(converted_path):
                file_path = converted_path
            else:
                raise Exception("Failed to convert .ppt to .pptx")
        except Exception as e:
            raise Exception(f"Old PowerPoint format (.ppt) detected but conversion failed: {str(e)}. Please convert to .pptx format first.")
    
    try:
        prs = Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        text = "\n".join(text_parts)
    except Exception as e:
        raise Exception(f"Failed to extract text from PowerPoint: {str(e)}")
    return text

def summarize_text(text, summary_length='medium', sentence_count=5):
    """Summarize text using extractive summarization"""
    try:
        from sumy.parsers.plaintext import PlaintextParser
        from sumy.nlp.tokenizers import Tokenizer
        from sumy.summarizers.lsa import LsaSummarizer
        from sumy.summarizers.text_rank import TextRankSummarizer
        from sumy.nlp.stemmers import Stemmer
        from sumy.utils import get_stop_words
        import nltk
        
        # Download required NLTK data if not already present
        try:
            nltk.data.find('tokenizers/punkt')
        except LookupError:
            nltk.download('punkt', quiet=True)
        
        # Determine sentence count based on summary length
        if summary_length == 'short':
            sentence_count = min(3, sentence_count)
        elif summary_length == 'medium':
            sentence_count = min(10, max(5, sentence_count))
        else:  # long
            sentence_count = min(20, max(10, sentence_count))
        
        # Parse text
        language = "english"
        parser = PlaintextParser.from_string(text, Tokenizer(language))
        stemmer = Stemmer(language)
        
        # Use TextRank summarizer (better quality than LSA)
        summarizer = TextRankSummarizer(stemmer)
        summarizer.stop_words = get_stop_words(language)
        
        # Generate summary
        summary_sentences = summarizer(parser.document, sentence_count)
        summary = " ".join([str(sentence) for sentence in summary_sentences])
        
        return summary, len(summary_sentences)
    
    except ImportError:
        # Fallback to simple sentence extraction if sumy is not available
        sentences = text.split('.')
        sentences = [s.strip() for s in sentences if s.strip()]
        summary_sentences = sentences[:sentence_count]
        summary = ". ".join(summary_sentences)
        if summary and not summary.endswith('.'):
            summary += "."
        return summary, len(summary_sentences)
    except Exception as e:
        raise Exception(f"Summarization failed: {str(e)}")

@app.route('/summarize/document', methods=['POST'])
def summarize_document():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        
        file = request.files['file']
        summary_length = request.form.get('summary_length', 'medium')
        sentence_count = int(request.form.get('sentence_count', 5))
        
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        # Validate file type
        file_ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
        if file_ext not in ['pdf', 'docx', 'pptx']:
            return jsonify({'error': 'Unsupported file type. Please upload PDF, DOCX, or PPTX files.'}), 400
        
        # Save uploaded file temporarily (use tempfile for Render compatibility)
        filename = secure_filename(file.filename)
        # Use tempfile for better compatibility with Render's ephemeral filesystem
        upload_path = os.path.join(tempfile.gettempdir(), f"upload_{datetime.now().strftime('%Y%m%d%H%M%S')}_{filename}")
        file.save(upload_path)
        
        try:
            # Extract text based on file type
            if file_ext == 'pdf':
                text = extract_text_from_pdf(upload_path)
            elif file_ext == 'docx':
                text = extract_text_from_word(upload_path)
            elif file_ext == 'pptx':
                text = extract_text_from_pptx(upload_path)
            else:
                os.remove(upload_path)
                return jsonify({'error': 'Unsupported file type'}), 400
            
            if not text.strip():
                os.remove(upload_path)
                return jsonify({'error': 'No text could be extracted from the document'}), 400
            
            # Count original sentences
            original_sentences = len([s for s in text.split('.') if s.strip()])
            
            # Generate summary
            summary, summary_sentence_count = summarize_text(text, summary_length, sentence_count)
            
            # Clean up uploaded file
            os.remove(upload_path)
            
            return jsonify({
                'summary': summary,
                'original_sentences': original_sentences,
                'summary_sentences': summary_sentence_count,
                'summary_length': summary_length
            })
        
        except Exception as e:
            if os.path.exists(upload_path):
                os.remove(upload_path)
            return jsonify({'error': str(e)}), 500
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

def generate_image_with_stable_diffusion(prompt, size='768x512', style='realistic', num_images=1):
    """Generate image using Stable Diffusion (if available)"""
    try:
        from diffusers import StableDiffusionPipeline
        import torch
        from PIL import Image
        import io
        import base64
        
        # Parse size
        width, height = map(int, size.split('x'))
        
        # Load model (this will download on first use)
        pipe = StableDiffusionPipeline.from_pretrained(
            "runwayml/stable-diffusion-v1-5",
            torch_dtype=torch.float16 if torch.cuda.is_available() else torch.float32
        )
        
        if torch.cuda.is_available():
            pipe = pipe.to("cuda")
        
        # Adjust prompt based on style
        style_prompts = {
            'realistic': f"photorealistic, high quality, {prompt}",
            'artistic': f"artistic, painting style, {prompt}",
            'cartoon': f"cartoon style, animated, {prompt}",
            'abstract': f"abstract art, modern, {prompt}"
        }
        enhanced_prompt = style_prompts.get(style, prompt)
        
        # Generate image
        images = pipe(enhanced_prompt, width=width, height=height, num_images_per_prompt=num_images).images
        
        # Convert to base64
        if num_images == 1:
            img = images[0]
            buffered = io.BytesIO()
            img.save(buffered, format="PNG")
            img_str = base64.b64encode(buffered.getvalue()).decode()
            return f"data:image/png;base64,{img_str}"
        else:
            # Return multiple images (for now, return first one)
            img = images[0]
            buffered = io.BytesIO()
            img.save(buffered, format="PNG")
            img_str = base64.b64encode(buffered.getvalue()).decode()
            return f"data:image/png;base64,{img_str}"
    
    except ImportError:
        # Fallback to placeholder if diffusers not available
        return generate_placeholder_image(prompt, size, style)
    except Exception as e:
        # Fallback on error
        print(f"Stable Diffusion error: {e}")
        return generate_placeholder_image(prompt, size, style)

def generate_placeholder_image(prompt, size='768x512', style='realistic'):
    """Generate a placeholder image using PIL (fallback)"""
    try:
        from PIL import Image, ImageDraw, ImageFont
        import io
        import base64
        import hashlib
        
        # Parse size
        width, height = map(int, size.split('x'))
        
        # Create image with gradient background
        img = Image.new('RGB', (width, height), color='#0f172a')
        draw = ImageDraw.Draw(img)
        
        # Create gradient effect
        for i in range(height):
            r = int(15 + (i / height) * 20)
            g = int(23 + (i / height) * 30)
            b = int(42 + (i / height) * 40)
            color = (r, g, b)
            draw.line([(0, i), (width, i)], fill=color)
        
        # Add text
        try:
            # Try to use a nicer font
            font_size = min(width // 20, 32)
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", font_size)
        except:
            try:
                font = ImageFont.load_default()
            except:
                font = None
        
        # Center text
        text = prompt[:50] + "..." if len(prompt) > 50 else prompt
        text_lines = []
        words = text.split()
        current_line = ""
        
        for word in words:
            test_line = current_line + (" " if current_line else "") + word
            if len(test_line) * 8 < width - 40:
                current_line = test_line
            else:
                if current_line:
                    text_lines.append(current_line)
                current_line = word
        if current_line:
            text_lines.append(current_line)
        
        # Draw text
        y_offset = (height - len(text_lines) * 30) // 2
        for i, line in enumerate(text_lines):
            bbox = draw.textbbox((0, 0), line, font=font) if font else (0, 0, len(line) * 8, 20)
            text_width = bbox[2] - bbox[0]
            x = (width - text_width) // 2
            y = y_offset + i * 30
            draw.text((x, y), line, fill='#60a5fa', font=font)
        
        # Add style indicator
        style_text = f"Style: {style.title()}"
        draw.text((20, 20), style_text, fill='#cbd5e1', font=font)
        
        # Convert to base64
        buffered = io.BytesIO()
        img.save(buffered, format="PNG")
        img_str = base64.b64encode(buffered.getvalue()).decode()
        return f"data:image/png;base64,{img_str}"
    
    except Exception as e:
        # Ultimate fallback - return error
        raise Exception(f"Image generation failed: {str(e)}")

@app.route('/generate/image', methods=['POST'])
def generate_image():
    try:
        data = request.get_json()
        prompt = data.get('prompt', '').strip()
        size = data.get('size', '768x512')
        style = data.get('style', 'realistic')
        num_images = int(data.get('num_images', 1))
        
        if not prompt:
            return jsonify({'error': 'Prompt is required'}), 400
        
        if len(prompt) < 10:
            return jsonify({'error': 'Prompt must be at least 10 characters long'}), 400
        
        # Validate size
        valid_sizes = ['512x512', '768x512', '512x768', '1024x1024']
        if size not in valid_sizes:
            size = '768x512'
        
        # Validate style
        valid_styles = ['realistic', 'artistic', 'cartoon', 'abstract']
        if style not in valid_styles:
            style = 'realistic'
        
        # Limit num_images
        num_images = min(max(1, num_images), 4)
        
        try:
            # Try Stable Diffusion first
            image_data = generate_image_with_stable_diffusion(prompt, size, style, num_images)
        except Exception as e:
            # Fallback to placeholder
            print(f"Image generation error: {e}")
            image_data = generate_placeholder_image(prompt, size, style)
        
        return jsonify({
            'image_data': image_data,
            'prompt': prompt,
            'size': size,
            'style': style
        })
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    # Get port from environment variable (Render provides this) or default to 8080
    port = int(os.environ.get('PORT', 8080))
    # Disable debug mode in production (Render sets this automatically)
    debug = os.environ.get('FLASK_ENV') == 'development' or os.environ.get('DEBUG', 'False').lower() == 'true'
    
    app.run(debug=debug, host='0.0.0.0', port=port)
